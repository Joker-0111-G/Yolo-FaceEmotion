"""
Yolo-FaceEmotion 答辩PPT生成脚本
生成一份美观大方、内容充实的答辩演示文稿
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
import os

# ==================== 色彩方案 ====================
# 主色调：深海蓝 + 活力橙
COLOR_PRIMARY = RGBColor(0x1A, 0x23, 0x3E)       # 深蓝黑 (背景/标题)
COLOR_SECONDARY = RGBColor(0x2D, 0x3A, 0x5C)     # 中蓝
COLOR_ACCENT = RGBColor(0xFF, 0x6B, 0x35)        # 活力橙 (强调)
COLOR_ACCENT2 = RGBColor(0x00, 0xB4, 0xD8)       # 青色 (次要强调)
COLOR_ACCENT3 = RGBColor(0x48, 0xCA, 0xE4)       # 浅青
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT_BG = RGBColor(0xF0, 0xF4, 0xF8)      # 浅灰蓝背景
COLOR_DARK_TEXT = RGBColor(0x1A, 0x23, 0x3E)
COLOR_GRAY_TEXT = RGBColor(0x6B, 0x72, 0x84)
COLOR_GREEN = RGBColor(0x2E, 0xCC, 0x71)          # 成功绿
COLOR_RED = RGBColor(0xE7, 0x4C, 0x3C)            # 警示红
COLOR_YELLOW = RGBColor(0xF3, 0x9C, 0x12)         # 警告黄
COLOR_PURPLE = RGBColor(0x9B, 0x59, 0xB6)         # 紫色

# 幻灯片尺寸 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


# ==================== 工具函数 ====================

def add_blank_slide():
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank layout
    return prs.slides.add_slide(layout)


def add_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, radius=None):
    """添加矩形形状"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    return shape


def add_circle(slide, left, top, size, fill_color=None):
    """添加圆形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=18, color=COLOR_DARK_TEXT,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei", anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    # 设置段落间距
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    return tf


def add_multiline_textbox(slide, left, top, width, height, lines, font_name="Microsoft YaHei"):
    """
    lines: list of (text, font_size, color, bold, alignment)
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        text, font_size, color, bold = line_data[:4]
        alignment = line_data[4] if len(line_data) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(4)
        p.space_before = Pt(2)
    return tf


def add_line(slide, x1, y1, x2, y2, color=COLOR_ACCENT, width=Pt(2)):
    """添加连接线"""
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_slide_number(slide, num):
    """在右下角添加页码"""
    add_textbox(slide, Inches(12.1), Inches(7.0), Inches(1.0), Inches(0.4),
                str(num), font_size=10, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.RIGHT)


def add_section_header(slide, title, subtitle="", slide_num=1):
    """现代风格的章节封面"""
    # 全屏深色背景
    add_bg(slide, COLOR_PRIMARY)

    # 左侧装饰条
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, fill_color=COLOR_ACCENT)

    # 顶部装饰线
    add_rect(slide, Inches(1.5), Inches(2.0), Inches(1.0), Inches(0.04), fill_color=COLOR_ACCENT)

    # 标题
    add_textbox(slide, Inches(1.5), Inches(2.3), Inches(10.5), Inches(1.2),
                title, font_size=44, color=COLOR_WHITE, bold=True)

    # 副标题
    if subtitle:
        add_textbox(slide, Inches(1.5), Inches(3.5), Inches(10.5), Inches(0.8),
                    subtitle, font_size=20, color=COLOR_GRAY_TEXT, bold=False)

    # 右下角小方块装饰
    add_rect(slide, Inches(12.0), Inches(6.5), Inches(0.6), Inches(0.6), fill_color=COLOR_ACCENT)
    add_rect(slide, Inches(11.2), Inches(6.5), Inches(0.6), Inches(0.6), fill_color=COLOR_ACCENT2)

    add_slide_number(slide, slide_num)


def add_content_slide_bg(slide, title, slide_num):
    """标准内容页背景"""
    add_bg(slide, COLOR_WHITE)

    # 顶部标题栏
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.2), fill_color=COLOR_PRIMARY)

    # 标题栏底部细线
    add_rect(slide, Inches(0), Inches(1.2), SLIDE_W, Inches(0.04), fill_color=COLOR_ACCENT)

    # 标题文字
    add_textbox(slide, Inches(0.8), Inches(0.15), Inches(11.0), Inches(0.9),
                title, font_size=30, color=COLOR_WHITE, bold=True)

    # 页码
    add_textbox(slide, Inches(12.1), Inches(7.05), Inches(1.0), Inches(0.35),
                str(slide_num), font_size=10, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.RIGHT)


def add_card(slide, left, top, width, height, title, content_lines, icon_text="",
             accent_color=COLOR_ACCENT, title_color=COLOR_PRIMARY):
    """添加卡片组件"""
    # 卡片背景
    card = add_rect(slide, left, top, width, height, fill_color=COLOR_LIGHT_BG, radius=Inches(0.15))

    # 顶部色条
    add_rect(slide, left + Inches(0.15), top + Inches(0.15), Inches(0.06), Inches(0.35),
             fill_color=accent_color)

    # 图标/标题
    title_text = f"{icon_text} {title}" if icon_text else title
    add_textbox(slide, left + Inches(0.4), top + Inches(0.12), width - Inches(0.6), Inches(0.45),
                title_text, font_size=18, color=title_color, bold=True)

    # 内容
    y_offset = Inches(0.65)
    for line in content_lines:
        add_textbox(slide, left + Inches(0.4), top + y_offset, width - Inches(0.6), Inches(0.3),
                    line, font_size=12, color=COLOR_GRAY_TEXT)
        y_offset += Inches(0.25)

    return card


# ==================== 第1页：封面 ====================
slide = add_blank_slide()
add_bg(slide, COLOR_PRIMARY)

# 大装饰圆
add_circle(slide, Inches(-2.0), Inches(-2.0), Inches(5.0), fill_color=RGBColor(0x22, 0x2E, 0x50))
add_circle(slide, Inches(10.5), Inches(5.0), Inches(4.0), fill_color=RGBColor(0x22, 0x2E, 0x50))

# 顶部细线
add_rect(slide, Inches(3.0), Inches(1.8), Inches(7.5), Inches(0.01), fill_color=RGBColor(0x3A, 0x4A, 0x70))

# 主标题
add_textbox(slide, Inches(3.0), Inches(2.0), Inches(7.5), Inches(1.5),
            "Yolo-FaceEmotion", font_size=52, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide, Inches(3.0), Inches(3.2), Inches(7.5), Inches(0.7),
            "基于YOLOv8和深度学习的人脸表情识别系统", font_size=24, color=COLOR_ACCENT3, bold=False, alignment=PP_ALIGN.CENTER)

# 分隔线
add_rect(slide, Inches(5.0), Inches(3.9), Inches(3.5), Inches(0.03), fill_color=COLOR_ACCENT)

# 课程信息
add_multiline_textbox(slide, Inches(3.0), Inches(4.3), Inches(7.5), Inches(1.8), [
    ("人工智能课程大作业 · 答辩汇报", 18, COLOR_GRAY_TEXT, False, PP_ALIGN.CENTER),
    ("", 8, COLOR_GRAY_TEXT, False, PP_ALIGN.CENTER),
    ("GitHub: github.com/Joker-0111-G/Yolo-FaceEmotion", 13, COLOR_ACCENT3, False, PP_ALIGN.CENTER),
])

# 底部
add_textbox(slide, Inches(3.0), Inches(6.5), Inches(7.5), Inches(0.5),
            "2025-2026 春季学期", font_size=13, color=COLOR_GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)


# ==================== 第2页：目录 ====================
slide = add_blank_slide()
add_section_header(slide, "汇报目录", "CONTENTS", 2)

toc_items = [
    ("01", "项目背景与意义", "人脸表情识别的应用价值与研究意义"),
    ("02", "数据集介绍", "FER2013数据集详情与数据分布"),
    ("03", "技术方案与模型架构", "YOLOv8s-cls + 迁移学习策略"),
    ("04", "模型训练与优化", "训练配置、超参数、训练过程"),
    ("05", "模型评估与结果分析", "准确率指标、类别分析、人类对比"),
    ("06", "实时推理系统", "系统流程、演示效果"),
    ("07", "挑战与展望", "技术挑战、未来优化方向"),
]

for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(2.0) + i * Inches(0.72)
    x = Inches(2.5)

    # 编号圆圈
    circle = add_circle(slide, x, y + Inches(0.05), Inches(0.42),
                        fill_color=COLOR_ACCENT if i == 0 else COLOR_SECONDARY)
    add_textbox(slide, x, y + Inches(0.08), Inches(0.42), Inches(0.42),
                num, font_size=16, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 标题
    add_textbox(slide, x + Inches(0.7), y, Inches(7.0), Inches(0.35),
                title, font_size=20, color=COLOR_WHITE, bold=True)

    # 描述
    add_textbox(slide, x + Inches(0.7), y + Inches(0.32), Inches(7.0), Inches(0.3),
                desc, font_size=12, color=COLOR_GRAY_TEXT)


# ==================== 第3页：项目背景与意义 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "项目背景与意义", 3)

# 左侧大卡片 - 什么是人脸表情识别
add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.5),
         "什么是人脸表情识别？", [
             "人脸表情识别 (Facial Expression Recognition, FER)",
             "是计算机视觉领域的重要研究方向。",
             "",
             "通过分析人脸图像中的面部特征（眉眼、嘴角、",
             "鼻翼等肌肉运动），自动判断人的情绪状态。",
             "",
             "7种基本表情：Anger · Disgust · Fear · Happy ·",
             "Neutral · Sad · Surprise (Ekman基本情绪理论)",
         ], icon_text="🎯", accent_color=COLOR_ACCENT)

# 右侧小卡片
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(1.0),
         "应用场景", [
             "🎓 智慧教育：学生课堂专注度与情绪分析",
             "🏥 心理健康：抑郁症等情绪障碍辅助筛查",
         ], icon_text="💡", accent_color=COLOR_ACCENT2)

add_card(slide, Inches(6.8), Inches(2.85), Inches(5.8), Inches(1.25),
         "", [
             "🚗 智能驾驶：驾驶员疲劳/情绪状态监测",
             "🤖 人机交互：情感计算与智能客服系统",
             "🛒 新零售：顾客满意度无感采集",
         ], icon_text="", accent_color=COLOR_ACCENT2)

# 底部 - 为什么选择这个项目
add_card(slide, Inches(0.6), Inches(4.4), Inches(12.0), Inches(2.5),
         "为什么选择这个项目？", [
             "1. 端到端实践：完整覆盖「数据 → 训练 → 评估 → 部署」的深度学习全流程",
             "2. 前沿技术栈：使用最新的YOLOv8架构 + PyTorch框架 + Ultralytics生态系统",
             "3. 迁移学习典范：ImageNet预训练 → FER2013微调，展示迁移学习在CV中的强大威力",
             "4. 实际可部署：不仅仅是模型训练，还实现了实时摄像头推理，可落地应用",
             "5. 超越人类基准：模型71.6%准确率 vs 人类平均~65%，证明了深度学习在该任务上的优势",
         ], icon_text="🚀", accent_color=COLOR_ACCENT)


# ==================== 第4页：数据集介绍 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "数据集介绍 — FER2013", 4)

# 左侧：数据集基本信息
add_card(slide, Inches(0.6), Inches(1.6), Inches(4.0), Inches(3.2),
         "FER2013 数据集概览", [
             "来源：Kaggle Facial Expression",
             "            Recognition 2013 Challenge",
             "",
             "总图片数：35,887 张",
             "图片尺寸：48×48 像素（灰度图）",
             "表情类别：7 类",
             "训练集：28,709 张 (80%)",
             "验证集：3,589 张 (10%)",
             "测试集：3,589 张 (10%)",
         ], icon_text="📊", accent_color=COLOR_ACCENT)

# 右侧：数据分布表格
add_card(slide, Inches(5.0), Inches(1.6), Inches(7.7), Inches(3.2),
         "各类别数据分布", [
             "类别     训练集   验证集   测试集    总计      占比",
             "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
             "Anger     3,995    467     491    4,953   13.8%",
             "Disgust     436     56      55      547    1.5%  ⚠️ 最少",
             "Fear      4,097    496     528    5,121   14.3%",
             "Happy     7,215    895     879    8,989   25.0%  ⭐ 最多",
             "Neutral   4,965    607     626    6,198   17.3%",
             "Sad       4,830    653     594    6,077   16.9%",
             "Surprise  3,171    415     416    4,002   11.2%",
         ], icon_text="📈", accent_color=COLOR_ACCENT2)

# 底部：关键发现
add_card(slide, Inches(0.6), Inches(5.1), Inches(12.0), Inches(1.8),
         "关键发现：严重的类别不平衡", [
             "😊 Happy 类有 7,215 张训练图片 → 模型最容易识别该类",
             "🤢 Disgust 类仅有 436 张训练图片 → 模型最难识别该类（样本量仅为 Happy 的 6%）",
             "📊 最大类/最小类比例 = 16.5:1，这是影响模型整体准确率的核心瓶颈之一",
             "💡 后续优化方向：对少数类（Disgust, Surprise）进行过采样或使用类别权重加权训练",
         ], icon_text="⚠️", accent_color=COLOR_YELLOW)


# ==================== 第5页：技术方案 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "技术方案与模型架构", 5)

# 左侧：YOLOv8介绍
add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(2.5),
         "YOLOv8s-cls 分类模型", [
             "YOLOv8 是 Ultralytics 公司开发的",
             "最新一代目标检测/分类/分割框架",
             "",
             "本项目使用 YOLOv8s-cls：",
             "• s = small (轻量版，~5M参数)",
             "• cls = 分类任务专用变体",
             "• 30层网络，12.5 GFLOPs计算量",
             "• 核心模块：Conv + C2f (跨阶段部分连接)",
             "• 输入尺寸：64×64 (FER2013原图48×48",
             "  上采样以保留更多面部特征)",
         ], icon_text="🧠", accent_color=COLOR_ACCENT)

# 右侧：迁移学习
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.5),
         "迁移学习策略", [
             "Transfer Learning 核心思想：",
             "将在大型数据集上训练的模型知识",
             "迁移到目标任务上",
             "",
             "本项目的迁移学习流程：",
             "1. 加载 yolov8s-cls.pt (ImageNet",
             "   预训练，1000类，百万级图片)",
             "2. 冻结156/158层预训练权重",
             "3. 替换最后的分类头 (1000类 → 7类)",
             "4. 在FER2013上微调200个epoch",
             "5. 仅训练约2小时即收敛",
         ], icon_text="🔄", accent_color=COLOR_ACCENT2)

# 底部：为什么选YOLO
add_card(slide, Inches(0.6), Inches(4.4), Inches(12.0), Inches(2.5),
         "技术选型理由", [
             "✅ 相比传统CNN（LeNet/AlexNet/VGG）: YOLOv8 架构更现代，C2f模块特征提取能力更强",
             "✅ 相比ViT (Vision Transformer): YOLO 在中小数据集上收敛更快，不需要海量预训练数据",
             "✅ 相比ResNet/EfficientNet: YOLO 内置RandAugment等增强策略，一个API调用即可启用",
             "✅ 生态成熟: Ultralytics框架提供完整的 train/val/predict/export 工具链，降低开发门槛",
             "✅ 推理速度: 0.3ms/张(GPU)，可满足实时视频流处理需求(30fps+)",
         ], icon_text="✅", accent_color=COLOR_GREEN)


# ==================== 第6页：系统架构 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "系统架构与数据流", 6)

# 流程图 - 用色块和箭头表示
flow_steps = [
    ("📷\n摄像头", COLOR_PRIMARY, Inches(0.5)),
    ("🔄\n灰度转换", COLOR_SECONDARY, Inches(2.0)),
    ("🔍\n人脸检测\nHaar Cascade", COLOR_ACCENT, Inches(3.5)),
    ("✂️\n人脸裁剪\n+15px padding", COLOR_ACCENT2, Inches(5.2)),
    ("🧠\nYOLO推理\n64×64 input", COLOR_PURPLE, Inches(6.9)),
    ("🏷️\n情绪标签\nTop-1 + 置信度", COLOR_GREEN, Inches(8.6)),
    ("🖼️\n画面渲染\n框+标签", COLOR_PRIMARY, Inches(10.3)),
]

for text, color, x in flow_steps:
    # 方块
    box = add_rect(slide, x, Inches(2.5), Inches(1.3), Inches(1.3), fill_color=color, radius=Inches(0.2))

    # 文字
    add_textbox(slide, x, Inches(2.6), Inches(1.3), Inches(1.1),
                text, font_size=13, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 箭头（除最后一个）
    if x < Inches(10.0):
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + Inches(1.35), Inches(2.95), Inches(0.35), Inches(0.4)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()

# 底部说明
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.5),
            "▲ 推理流程图：从摄像头采集到最终画面渲染的完整数据通路",
            font_size=11, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# 下方：两阶段架构
add_card(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(2.2),
         "阶段一：人脸检测（传统CV）", [
             "• 算法：OpenCV Haar Cascade 级联分类器",
             "• 原理：Haar-like特征 + AdaBoost级联",
             "• 参数：scaleFactor=1.1, minNeighbors=5",
             "• 最小人脸尺寸：60×60 像素",
             "• 优势：极快，CPU即可实时运行",
             "• 局限：对侧脸、遮挡、极端光照敏感",
         ], icon_text="🔍", accent_color=COLOR_ACCENT)

add_card(slide, Inches(6.8), Inches(4.8), Inches(5.8), Inches(2.2),
         "阶段二：表情分类（深度学习）", [
             "• 模型：YOLOv8s-cls（~5M参数）",
             "• 输入：裁剪后的人脸区域 → 64×64",
             "• 输出：7类表情的概率分布 + Top-1标签",
             "• 推理速度：0.3ms/张（GPU）",
             "• 预处理：0.0ms（YOLO内置标准化）",
             "• 后处理：0.0ms（直接取softmax最大值）",
         ], icon_text="🧠", accent_color=COLOR_PURPLE)


# ==================== 第7页：模型训练 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "模型训练与优化", 7)

# 训练配置表
add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(3.5),
         "训练超参数配置", [
             "模型基座：yolov8s-cls.pt (ImageNet预训练)",
             "输入尺寸：64 × 64 像素",
             "训练轮数：200 epochs",
             "批次大小：32（兼顾显存与收敛速度）",
             "优化器：MuSGD (Ultralytics自动选择)",
             "初始学习率：0.01 (warm-up 3 epochs)",
             "学习率调度：cos_lr (余弦退火)",
             "数据增强：RandAugment（自动搜索增强策略）",
             "Dropout：0.0（使用RandAugment替代正则化）",
             "训练时间：2.102 小时 (RTX 3050 4GB)",
             "GPU显存占用：~1.0 GB（峰值~2GB）",
             "训练集：28,709张 / 验证集：3,589张",
         ], icon_text="⚙️", accent_color=COLOR_ACCENT)

# 训练过程曲线（用文字描述 + 数据点）
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.5),
         "训练过程（损失下降曲线）", [
             "Epoch   GPU_mem   loss   变化趋势",
             "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
             "  1      1.0G    1.1482  ████████████████",
             " 50      1.0G    0.5450  ████████ (↓52.5%)",
             "100      1.0G    0.4115  ██████   (↓64.2%)",
             "150      1.0G    0.3523  █████    (↓69.3%)",
             "200      1.0G    0.3143  ████     (↓72.6%)",
             "",
             "▸ 前50轮：损失快速下降（模型快速学习）",
             "▸ 50-150轮：稳定下降（精细调优阶段）",
             "▸ 150-200轮：收敛缓慢（接近最优）",
             "▸ 最优模型出现在 ~Epoch 110-120",
         ], icon_text="📉", accent_color=COLOR_ACCENT2)

# 底部策略
add_card(slide, Inches(0.6), Inches(5.35), Inches(12.0), Inches(1.6),
         "训练优化策略", [
             "🔹 RandAugment：自动搜索最优数据增强组合（random erasing, rotation, color jitter等），替代手工设计增强 → 有效防止过拟合",
             "🔹 MuSGD (Momentum SGD)：带动量的随机梯度下降，momentum=0.937 → 加速收敛、减少震荡",
             "🔹 Warm-up (3 epochs)：学习率从0.01×0.1线性升至0.01 → 避免训练初期不稳定",
             "🔹 Early Stopping (patience=100)：监控验证集准确率，超过100轮不提升则自动停止 → 防止过拟合",
         ], icon_text="🔧", accent_color=COLOR_YELLOW)


# ==================== 第8页：模型评估 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "模型评估与结果分析", 8)

# 核心指标大卡
metrics_data = [
    ("71.61%", "Top-1 准确率", COLOR_ACCENT),
    ("99.20%", "Top-5 准确率", COLOR_GREEN),
    ("0.3ms", "单张推理速度", COLOR_ACCENT2),
    ("2.1h", "训练总时长", COLOR_PURPLE),
]

for i, (value, label, color) in enumerate(metrics_data):
    x = Inches(0.8) + i * Inches(3.1)
    # 背景圆角矩形
    add_rect(slide, x, Inches(1.6), Inches(2.7), Inches(1.6), fill_color=COLOR_LIGHT_BG, radius=Inches(0.15))
    # 数值
    add_textbox(slide, x, Inches(1.8), Inches(2.7), Inches(0.7),
                value, font_size=36, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    # 标签
    add_textbox(slide, x, Inches(2.6), Inches(2.7), Inches(0.4),
                label, font_size=14, color=COLOR_GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# 左下：准确率提升过程
add_card(slide, Inches(0.6), Inches(3.5), Inches(6.0), Inches(2.2),
         "准确率提升路径", [
             "Epoch   验证集Top-1   提升",
             "━━━━━━━━━━━━━━━━━━━━━━━━━━",
             "初始(预训练)  ~55%      -",
             " 50 epochs    ~64%      +9%",
             "100 epochs    ~67%      +3%",
             "150 epochs    ~69%      +2%",
             "200 epochs    ~70.3%    +1.3%",
             "",
             "最终测试集 Top-1: 71.61%（泛化良好）",
         ], icon_text="📈")

# 右下：关键发现
add_card(slide, Inches(6.8), Inches(3.5), Inches(5.8), Inches(2.2),
         "评估关键发现", [
             "✅ Top-1 71.61% → 在7分类任务上表现良好",
             "✅ Top-5 99.2% → 模型几乎从不把表情",
             "     误判为完全不相关的类别",
             "⚠️ 验证集约70.3% vs 测试集71.61%",
             "     → 测试集表现更好，无过拟合",
             "🔥 推理仅需0.3ms → 可支持",
             "     3000+ fps 的实时处理速度",
         ], icon_text="🔍", accent_color=COLOR_ACCENT2)

# 底部混淆分析
add_card(slide, Inches(0.6), Inches(5.95), Inches(12.0), Inches(1.0),
         "主要误判模式分析", [
             "🔸 Fear ↔ Surprise 混淆：两种表情面部特征相似（眉眼上扬、嘴巴张开），是最容易混淆的类别对",
             "🔸 Sad ↔ Neutral 混淆：轻度悲伤与无表情难以区分，尤其是灰度低分辨率图片中",
             "🔸 Anger ↔ Disgust 混淆：两种负面情绪都有皱眉、鼻翼收缩等特征，且Disgust训练样本极少",
         ], icon_text="🔄", accent_color=COLOR_RED)


# ==================== 第9页：各类别表现 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "各类别表现分析", 9)

# 类别表现卡片
class_perf = [
    ("😊 Happy", "⭐ 最高", "样本最多 (7,215)，特征明显易识别", COLOR_GREEN),
    ("😐 Neutral", "⭐ 较高", "样本充足 (4,965)，面部特征简单", COLOR_GREEN),
    ("😠 Anger", "✅ 良好", "样本适中 (3,995)，特征清晰", COLOR_ACCENT2),
    ("😲 Surprise", "✅ 良好", "特征鲜明，但样本偏少 (3,171)", COLOR_ACCENT2),
    ("😢 Sad", "⚠️ 中等", "轻度悲伤与Neutral易混淆", COLOR_YELLOW),
    ("😨 Fear", "⚠️ 中等", "常与Surprise混淆，特征相似", COLOR_YELLOW),
    ("🤢 Disgust", "🔻 最低", "仅436张训练图（最少），严重欠拟合", COLOR_RED),
]

for i, (emoji_label, perf, reason, color) in enumerate(class_perf):
    x = Inches(0.6) + (i % 4) * Inches(3.1)
    y = Inches(1.6) + (i // 4) * Inches(2.7)
    w = Inches(2.8)
    h = Inches(2.4)

    add_rect(slide, x, y, w, h, fill_color=COLOR_LIGHT_BG, radius=Inches(0.12))

    # 顶部色条
    add_rect(slide, x + Inches(0.15), y + Inches(0.12), Inches(0.05), Inches(0.35), fill_color=color)

    # 表情标签
    add_textbox(slide, x + Inches(0.35), y + Inches(0.1), w - Inches(0.5), Inches(0.5),
                emoji_label, font_size=20, color=COLOR_PRIMARY, bold=True)

    # 表现评级
    add_textbox(slide, x + Inches(0.35), y + Inches(0.55), w - Inches(0.5), Inches(0.35),
                perf, font_size=13, color=color, bold=True)

    # 原因
    add_textbox(slide, x + Inches(0.35), y + Inches(1.0), w - Inches(0.5), Inches(1.2),
                reason, font_size=11, color=COLOR_GRAY_TEXT)

# 底部总结
add_textbox(slide, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.5),
            "💡 类别准确率与训练样本数量呈正相关 → 解决类别不平衡是提升模型性能的关键方向",
            font_size=13, color=COLOR_DARK_TEXT, bold=True)


# ==================== 第10页：与人类基准对比 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "与人类基准对比", 10)

# 中心对比图
# 人类
human_card = add_rect(slide, Inches(1.5), Inches(2.0), Inches(4.0), Inches(3.5),
                      fill_color=COLOR_LIGHT_BG, radius=Inches(0.2))
add_textbox(slide, Inches(1.5), Inches(2.2), Inches(4.0), Inches(0.6),
            "👤 人类平均水平", font_size=22, color=COLOR_PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(3.0), Inches(4.0), Inches(1.2),
            "~65%", font_size=60, color=COLOR_RED, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(4.0), Inches(0.8),
            "Top-1 准确率\n(FER2013 数据集)", font_size=13, color=COLOR_GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# VS
add_textbox(slide, Inches(5.5), Inches(3.2), Inches(2.5), Inches(1.0),
            "VS", font_size=48, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

# 模型
model_card = add_rect(slide, Inches(8.0), Inches(2.0), Inches(4.0), Inches(3.5),
                      fill_color=COLOR_PRIMARY, radius=Inches(0.2))
add_textbox(slide, Inches(8.0), Inches(2.2), Inches(4.0), Inches(0.6),
            "🤖 Yolo-FaceEmotion", font_size=22, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.0), Inches(3.0), Inches(4.0), Inches(1.2),
            "~71.6%", font_size=60, color=COLOR_GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(8.0), Inches(4.2), Inches(4.0), Inches(0.8),
            "Top-1 准确率\n(FER2013 测试集)", font_size=13, color=COLOR_ACCENT3, bold=False, alignment=PP_ALIGN.CENTER)

# 提升箭头
add_textbox(slide, Inches(3.0), Inches(5.8), Inches(7.5), Inches(0.5),
            "↑ 模型超越人类约 6.6 个百分点（相对提升 10.2%）",
            font_size=18, color=COLOR_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# 底部说明
add_textbox(slide, Inches(1.0), Inches(6.5), Inches(11.5), Inches(0.8),
            "注：人类在FER2013上的平均准确率约为 65% ± 5%（来源：FER2013竞赛报告及后续研究）。"
            "模型通过学习数万张表情图片，掌握了细微的面部特征差异，超越了未经训练的人类观察者。",
            font_size=10, color=COLOR_GRAY_TEXT, alignment=PP_ALIGN.CENTER)


# ==================== 第11页：实时推理系统 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "实时推理系统", 11)

# 左侧：系统特点
add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(3.0),
         "实时推理系统 (main.py)", [
             "🎥 调用本地摄像头 (cv2.VideoCapture(0))",
             "🔄 每帧转为灰度图 → 加速人脸检测",
             "🔍 Haar Cascade 检测画面中所有人脸",
             "✂️ 每张人脸向外扩展15px → 截取",
             "🧠 YOLO前向推理 → 获取情绪分类",
             "🖼️ 绿色框 + 标签 + 置信度叠加渲染",
             "⌨️ 按 'q' 键退出，自动释放资源",
         ], icon_text="🎥", accent_color=COLOR_ACCENT)

# 右侧：演示效果示意
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(3.0),
         "显示效果示意", [
             "┌──────────────────────────────┐",
             "│                              │",
             "│     ┌──────────┐             │",
             "│     │  😊      │             │",
             "│     │ Happy    │ ← 绿色框    │",
             "│     │ (98.3%)  │ ← 置信度    │",
             "│     └──────────┘             │",
             "│                              │",
             "│   窗口标题:                   │",
             "│   YOLO Emotion Recognition   │",
             "└──────────────────────────────┘",
         ], icon_text="🖥️", accent_color=COLOR_ACCENT2)

# 底部：技术亮点
add_card(slide, Inches(0.6), Inches(4.9), Inches(5.8), Inches(2.0),
         "工程亮点", [
             "✅ 灰度转换加速：人脸检测在灰度图进行，",
             "    减少2/3计算量",
             "✅ verbose=False：抑制每帧控制台输出，",
             "    避免刷屏影响性能",
             "✅ padding=15px：防止人脸边缘特征丢失",
             "✅ 边界保护：max(0, y-padding) 防止越界",
         ], icon_text="⚡", accent_color=COLOR_GREEN)

add_card(slide, Inches(6.8), Inches(4.9), Inches(5.8), Inches(2.0),
         "性能指标", [
             "📹 帧处理速度：30+ fps（取决于摄像头帧率）",
             "⏱️ 人脸检测：~5ms（Haar Cascade, CPU）",
             "🚀 YOLO推理：0.3ms/张（GPU, RTX 3050）",
             "📊 总延迟：~10ms/帧 → 完全满足实时需求",
             "💾 GPU显存占用：< 1GB（推理模式）",
             "🖥️ CPU占用：低（主要负载在GPU推理）",
         ], icon_text="📊", accent_color=COLOR_PURPLE)


# ==================== 第12页：技术挑战与解决方案 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "技术挑战与解决方案", 12)

challenges = [
    ("挑战一：类别严重不平衡",
     "Disgust类仅436张训练图，Happy类7,215张\n→ 最大/最小类比例 = 16.5:1",
     "✅ 使用RandAugment对所有类别统一增强\n🔜 未来：对少数类过采样(SMOTE)或使用Focal Loss",
     COLOR_RED, COLOR_GREEN),

    ("挑战二：低分辨率图片",
     "FER2013原图仅48×48像素灰度图\n→ 面部细节有限，模型难以捕捉微妙差异",
     "✅ 上采样至64×64，保留更多特征信息\n✅ YOLO内部特征金字塔多尺度提取特征",
     COLOR_RED, COLOR_GREEN),

    ("挑战三：Haar Cascade人脸检测局限",
     "对侧脸、遮挡、极端光照场景检测率低\n→ 人脸检测失败 → 整个系统无法工作",
     "✅ 本项目当前使用Haar（速度优先）\n🔜 未来：替换为YOLO-face检测模型",
     COLOR_RED, COLOR_GREEN),

    ("挑战四：表情模糊性与主观性",
     "Fear/Surprise、Sad/Neutral 易混淆\n→ 即使是人类标注者也存在分歧",
     "✅ Top-5准确率99.2%：模型知道\"近似表情\"\n✅ Top-1 71.6%已超过人类平均水平",
     COLOR_RED, COLOR_GREEN),
]

for i, (title, desc, solution, c1, c2) in enumerate(challenges):
    y = Inches(1.6) + i * Inches(1.4)
    x = Inches(0.6)

    # 标题
    add_textbox(slide, x, y, Inches(5.0), Inches(0.35),
                title, font_size=16, color=c1, bold=True)

    # 描述
    add_textbox(slide, x + Inches(5.2), y, Inches(3.5), Inches(0.7),
                desc, font_size=10, color=COLOR_DARK_TEXT)

    # 解决方案
    add_textbox(slide, x + Inches(8.8), y, Inches(3.8), Inches(0.7),
                solution, font_size=10, color=c2)

    # 分隔线
    if i < len(challenges) - 1:
        add_line(slide, x, y + Inches(1.2), x + Inches(12.0), y + Inches(1.2),
                color=RGBColor(0xE0, 0xE4, 0xE8), width=Pt(0.5))


# ==================== 第13页：未来优化方向 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "未来优化方向", 13)

future_items = [
    ("🔬", "数据层面优化", [
        "对少数类(Disgust)过采样至与Happy相当",
        "引入更多真实场景数据集(RAF-DB, AffectNet)",
        "尝试GAN生成少数类表情图片",
        "使用更强的数据增强：CutMix + MixUp",
    ], COLOR_ACCENT),

    ("🏗️", "模型架构升级", [
        "尝试YOLOv8m-cls / YOLOv8l-cls更大模型",
        "对比ViT/EfficientNet等架构在该任务的表现",
        "引入注意力机制增强关键面部区域权重",
        "探索多任务学习：同时预测表情+关键点",
    ], COLOR_PURPLE),

    ("🔍", "人脸检测升级", [
        "Haar Cascade → YOLO人脸检测模型",
        "支持多角度人脸检测(侧脸、低头、仰头)",
        "添加人脸对齐预处理(MTCNN/RetinaFace)",
        "引入人脸质量评估过滤低质量输入",
    ], COLOR_ACCENT2),

    ("📱", "工程部署优化", [
        "FP16/INT8量化 → 模型体积缩小75%",
        "导出ONNX/TensorRT格式加速推理",
        "OpenVINO部署到边缘设备(树莓派/Jetson)",
        "Flask/Streamlit搭建Web交互界面",
    ], COLOR_GREEN),
]

for i, (icon, title, items, color) in enumerate(future_items):
    x = Inches(0.6) + i * Inches(3.1)
    y = Inches(1.6)

    add_rect(slide, x, y, Inches(2.85), Inches(5.0), fill_color=COLOR_LIGHT_BG, radius=Inches(0.15))

    # 图标
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.5), Inches(0.5),
                icon, font_size=24, alignment=PP_ALIGN.CENTER)

    # 标题
    add_textbox(slide, x + Inches(0.7), y + Inches(0.15), Inches(1.9), Inches(0.5),
                title, font_size=17, color=color, bold=True)

    # 色条
    add_rect(slide, x + Inches(0.15), y + Inches(0.75), Inches(1.0), Inches(0.03), fill_color=color)

    # 内容
    for j, item in enumerate(items):
        add_textbox(slide, x + Inches(0.2), y + Inches(1.0) + j * Inches(0.55),
                    Inches(2.4), Inches(0.6),
                    f"• {item}", font_size=11, color=COLOR_DARK_TEXT)


# ==================== 第14页：项目总结 ====================
slide = add_blank_slide()
add_content_slide_bg(slide, "项目总结与收获", 14)

# 左侧：项目成果
add_card(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(4.8),
         "项目成果", [
             "✅ 成功构建了端到端的人脸表情识别系统",
             "",
             "✅ 掌握了YOLOv8分类模型的训练流程",
             "   数据准备 → 超参数配置 → 训练监控 → 模型保存",
             "",
             "✅ 理解了迁移学习的核心思想和实践方法",
             "   ImageNet预训练 → FER2013微调 → 精度提升",
             "",
             "✅ 实现了实时摄像头推理的完整工程链路",
             "   图像采集 → 人脸检测 → 表情分类 → 结果渲染",
             "",
             "✅ 模型Top-1准确率71.61%，超越人类基准",
             "   达到了课程项目的预期目标",
             "",
             "✅ 深入理解了数据集质量对深度学习的影响",
             "   类别不平衡 → 少数类准确率低 → 优化方向",
         ], icon_text="🏆", accent_color=COLOR_ACCENT)

# 右侧：技能收获
add_card(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(4.8),
         "技能与知识收获", [
             "🔹 PyTorch & Ultralytics 框架实战经验",
             "",
             "🔹 YOLO系列模型架构理解(Conv + C2f + Head)",
             "",
             "🔹 计算机视觉全流程：",
             "   分类/检测/分割任务的异同与选择",
             "",
             "🔹 OpenCV图像处理：",
             "   颜色空间转换、人脸检测、图像渲染",
             "",
             "🔹 深度学习训练技巧：",
             "   RandAugment、Warm-up、Cosine LR、Early Stop",
             "",
             "🔹 模型评估方法论：",
             "   Top-1/Top-5、混淆矩阵、类别分析",
             "",
             "🔹 完整项目的工程化思维",
         ], icon_text="📚", accent_color=COLOR_ACCENT2)


# ==================== 第15页：致谢 & Q&A ====================
slide = add_blank_slide()
add_bg(slide, COLOR_PRIMARY)

# 装饰
add_circle(slide, Inches(-1.5), Inches(-1.5), Inches(4.0), fill_color=RGBColor(0x22, 0x2E, 0x50))
add_circle(slide, Inches(11.0), Inches(5.5), Inches(3.5), fill_color=RGBColor(0x22, 0x2E, 0x50))

# 大标题
add_textbox(slide, Inches(2.0), Inches(1.5), Inches(9.5), Inches(1.0),
            "感谢聆听", font_size=48, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_textbox(slide, Inches(2.0), Inches(2.5), Inches(9.5), Inches(0.6),
            "THANK YOU", font_size=28, color=COLOR_GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# 分隔线
add_rect(slide, Inches(5.0), Inches(3.2), Inches(3.5), Inches(0.03), fill_color=COLOR_ACCENT)

# 项目信息
add_multiline_textbox(slide, Inches(2.0), Inches(3.5), Inches(9.5), Inches(2.0), [
    ("Yolo-FaceEmotion", 24, COLOR_ACCENT3, True, PP_ALIGN.CENTER),
    ("基于YOLOv8s-cls和FER2013的实时人脸表情识别系统", 16, COLOR_GRAY_TEXT, False, PP_ALIGN.CENTER),
    ("", 10, COLOR_WHITE, False, PP_ALIGN.CENTER),
    ("GitHub: github.com/Joker-0111-G/Yolo-FaceEmotion", 13, COLOR_ACCENT3, False, PP_ALIGN.CENTER),
])

# Q&A
add_textbox(slide, Inches(2.0), Inches(5.5), Inches(9.5), Inches(0.6),
            "Q & A", font_size=32, color=COLOR_ACCENT, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(2.0), Inches(6.1), Inches(9.5), Inches(0.5),
            "欢迎提问与交流", font_size=16, color=COLOR_GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# 底部
add_textbox(slide, Inches(2.0), Inches(6.8), Inches(9.5), Inches(0.4),
            "人工智能课程大作业 · 2025-2026春季学期", font_size=11, color=COLOR_GRAY_TEXT, bold=False, alignment=PP_ALIGN.CENTER)

# ==================== 保存 ====================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Yolo-FaceEmotion_答辩PPT.pptx")
prs.save(output_path)
print(f"[OK] PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
